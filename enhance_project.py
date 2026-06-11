from __future__ import annotations

import json
from pathlib import Path


ROOT = Path(__file__).resolve().parent
APP = ROOT / "app" / "streamlit_app.py"
REPORT = ROOT / "reports" / "Scam_Detection_Report_2page.md"
FULL_REPORT = ROOT / "reports" / "Scam_Detection_Report.md"
NOTEBOOK = ROOT / "Scam_or_Not_Scam_IR_Project.ipynb"
SLIDES = ROOT / "slides" / "Scam_Detection_Presentation.pptx"


LITERATURE = [
    {
        "author": "T. A. Almeida, J. M. G. Hidalgo, and A. Yamakami",
        "title": "A Contribution to the Study of SMS Spam Filtering: New Collection and Results",
        "methodology": "Created a public SMS spam corpus and benchmarked traditional text-classification methods for ham/spam detection.",
        "insight": "Short-message classification can work well with carefully cleaned text, word dictionaries, and supervised classifiers.",
        "gap": "The dataset is older and mostly spam/ham, so modern scam and smishing messages require newer, domain-specific validation.",
        "url": "https://www.dt.fee.unicamp.br/~tiago/smsspamcollection/",
    },
    {
        "author": "N. Al Moubayed, T. Breckon, P. Matthews, and A. S. McGough",
        "title": "SMS Spam Filtering using Probabilistic Topic Modelling and Stacked Denoising Autoencoder",
        "methodology": "Combined topic modelling with stacked denoising autoencoders to classify SMS spam with limited manual feature engineering.",
        "insight": "Topic representations can improve interpretability, while learned representations can achieve strong spam-filtering accuracy.",
        "gap": "Deep models are less transparent for beginner users and may be unnecessary when a small labelled corpus is linearly separable.",
        "url": "https://arxiv.org/abs/1606.05554",
    },
    {
        "author": "Y. Li, R. Zhang, W. Rong, and X. Mi",
        "title": "SpamDam: Towards Privacy-Preserving and Adversary-Resistant SMS Spam Detection",
        "methodology": "Built a large SMS spam collection pipeline, studied campaign patterns, and evaluated centralized and federated spam detectors.",
        "insight": "Real-world SMS spam changes over time and privacy-preserving learning is important for deployment.",
        "gap": "Portfolio projects should add drift monitoring, feedback loops, and privacy controls before use with personal messages.",
        "url": "https://arxiv.org/abs/2404.09481",
    },
    {
        "author": "M. Salman, M. Ikram, N. Basta, and M. A. Kaafar",
        "title": "SpaLLM-Guard: Pairing SMS Spam Detection Using Open-source and Commercial LLMs",
        "methodology": "Compared zero-shot, few-shot, fine-tuned, and chain-of-thought LLM strategies for SMS spam detection.",
        "insight": "Fine-tuned language models can be robust, but zero-shot prompting alone is unreliable for production spam detection.",
        "gap": "LLMs are powerful but heavier than classical models; small TF-IDF/logistic models remain attractive for transparent, low-cost deployments.",
        "url": "https://arxiv.org/abs/2501.04985",
    },
    {
        "author": "D. Goel, H. Ahmad, A. K. Jain, and N. K. Goel",
        "title": "Machine Learning Driven Smishing Detection Framework for Mobile Security",
        "methodology": "Used content-based smishing detection with text normalization and machine-learning classifiers.",
        "insight": "Normalizing slang, abbreviations, and short forms improves mobile-message threat detection.",
        "gap": "The current project should later add Khmer-English normalization and multilingual scam examples.",
        "url": "https://arxiv.org/abs/2412.09641",
    },
]


def md_table(rows: list[dict]) -> str:
    header = "| Author | Title | Methodology | Key Insight | Research Gap |\n|---|---|---|---|---|"
    body = "\n".join(
        f"| {r['author']} | [{r['title']}]({r['url']}) | {r['methodology']} | {r['insight']} | {r['gap']} |"
        for r in rows
    )
    return header + "\n" + body


def enhance_report() -> None:
    metrics = json.loads((ROOT / "reports" / "metrics.json").read_text(encoding="utf-8"))
    comparison = (ROOT / "reports" / "model_comparison.csv").read_text(encoding="utf-8").strip().splitlines()
    comp_rows = []
    for line in comparison[1:]:
        model, accuracy, precision, recall, f1, roc_auc = line.split(",")
        comp_rows.append(f"| {model} | {float(accuracy):.4f} | {float(precision):.4f} | {float(recall):.4f} | {float(f1):.4f} | {float(roc_auc):.4f} |")
    comp_table = "\n".join(comp_rows)
    lit_table = md_table(LITERATURE[:4])

    report = f"""# Scam or Not Scam Text Classification Report

**Student:** [Your Name]  
**Student ID:** [Your Student ID]  
**Course:** Information Web Retrieval Analysis / Natural Language Processing  
**Project:** Mini Project 2 - Binary Text Classification  

## 1. Objective

This project adapts the assigned positive/negative review classification workflow to the supplied `safe` and `scam` message corpus. The goal is to classify a new message as **safe** or **scam** using an Information Retrieval pipeline: text preprocessing, dictionary construction, feature extraction, vector-space representation, supervised classification, evaluation, and a deployable Streamlit prototype.

## 2. Literature Review

{lit_table}

The literature shows that short-message scam detection benefits from a balance between **interpretable engineered features** and **statistical text representations**. Therefore, this project uses both assignment-required features and additional real-world scam indicators.

## 3. Dataset and Preprocessing

The dataset contains 40,000 labelled messages: 20,000 safe and 20,000 scam. Following the project instruction, the top 80% of each class is used for training and the final 20% for testing. Preprocessing includes lowercasing, regex tokenization, URL/phone/money detection, punctuation cues, pronoun matching, and dictionary matching against safe/scam word lists.

## 4. Features

Required assignment features:

- Count of safe/positive words.
- Count of scam/negative words.
- Binary indicator for the word `no`.
- Count of first/second person pronouns: `I`, `me`, `my`, `you`, `your`.
- Binary indicator for exclamation mark `!`.
- Log length of the message.

Additional professional features:

- URL, phone-number, and money/prize indicators.
- Urgency and action-word counts.
- Digit count, uppercase ratio, average word length, question count.
- Safe/scam word density.
- TF-IDF unigram and bigram representation for vector-space modeling.

## 5. Model Architecture and Experiments

Three models were compared:

| Model | Accuracy | Precision | Recall | F1 | ROC-AUC |
|---|---:|---:|---:|---:|---:|
{comp_table}

The selected model is **{metrics['model_type']}** because it achieved the highest test performance and remains transparent enough for a portfolio-ready safety application.

## 6. Results and Discussion

The final model achieved **{metrics['accuracy']:.4f} accuracy**, **{metrics['precision']:.4f} precision**, **{metrics['recall']:.4f} recall**, and **{metrics['f1']:.4f} F1-score**. The confusion matrix is `{metrics['confusion_matrix']}` in the order `[[safe-safe, safe-scam], [scam-safe, scam-scam]]`.

The perfect result suggests that the provided corpus is highly separable. This is useful for demonstrating the project workflow, but a real deployment should be tested on newer, messier, multilingual scam messages. For real users, recall is especially important because a missed scam can cause financial or privacy harm.

## 7. Demo and Deployment

The saved pipeline is connected to a Streamlit app. Users can paste a suspicious SMS, email, or chat message and receive a scam probability, warning label, risk signals, and safety advice. The app also includes a term explorer to search corpus vocabulary and inspect common safe/scam indicators.

## 8. Future Work

Future improvements include Khmer-English normalization, highlighted suspicious phrases, user feedback collection, drift monitoring, privacy-first deployment, and comparison with transformer embeddings or fine-tuned language models.
"""
    REPORT.write_text(report, encoding="utf-8")

    full = FULL_REPORT.read_text(encoding="utf-8")
    marker = "## Problem Statement"
    lit_section = "## Literature Review\n\n" + md_table(LITERATURE) + "\n\n"
    if "## Literature Review" not in full and marker in full:
        full = full.replace(marker, lit_section + marker)
    FULL_REPORT.write_text(full, encoding="utf-8")


def cell(cell_type: str, source: str) -> dict:
    data = {"cell_type": cell_type, "metadata": {}, "source": source.splitlines(True)}
    if cell_type == "code":
        data.update({"execution_count": None, "outputs": []})
    return data


def enhance_notebook() -> None:
    nb = json.loads(NOTEBOOK.read_text(encoding="utf-8"))
    intro = [
        cell("markdown", "# Scam or Not Scam Detection\n\n**Student:** [Your Name]  \n**Student ID:** [Your Student ID]  \n**Course:** Information Web Retrieval Analysis / NLP  \n\nThis notebook follows the Mini Project 2 instruction and extends it into a professional portfolio-ready scam-message detection workflow."),
        cell("markdown", "## Project Roadmap\n\n1. Understand the problem and assignment features.\n2. Load safe/scam text files.\n3. Clean and preprocess text.\n4. Perform EDA with tables and graphs.\n5. Extract required and additional features.\n6. Train/test split using top 80% per class.\n7. Compare multiple models.\n8. Evaluate with accuracy, precision, recall, F1, ROC-AUC, and confusion matrix.\n9. Run a demo prediction.\n10. Save model artifacts for Streamlit deployment."),
        cell("markdown", "## Literature Review Summary\n\nThe report includes four related studies on SMS spam and smishing detection. The main lesson is that classical ML models remain useful when they combine transparent features with strong text representations, while real deployment requires newer data, multilingual normalization, privacy protection, and drift monitoring."),
    ]
    if "Project Roadmap" not in "".join("".join(c.get("source", [])) for c in nb["cells"][:5]):
        nb["cells"] = intro + nb["cells"][1:]

    additions = [
        cell("markdown", "## Professional EDA: Label Balance and Message Length\n\nThe following plots make the dataset easier to understand before modeling. Good EDA should show class balance, text length distribution, and feature differences between safe and scam messages."),
        cell("code", "import seaborn as sns\nsns.set_theme(style='whitegrid', palette='Set2')\nfig, axes = plt.subplots(1, 2, figsize=(13, 4))\nsns.countplot(data=df_model, x='label', ax=axes[0])\naxes[0].set_title('Class Balance')\naxes[0].set_xlabel('Message Label')\naxes[0].set_ylabel('Count')\nsns.histplot(data=df_model, x='word_count', hue='label', bins=40, kde=True, ax=axes[1])\naxes[1].set_title('Message Length Distribution')\naxes[1].set_xlabel('Token Count')\nplt.tight_layout()"),
        cell("markdown", "## Required Feature Audit\n\nThis table confirms that every assignment-required feature is present in the dataset."),
        cell("code", "required_features = ['safe_word_hits', 'scam_word_hits', 'contains_no', 'first_second_pronoun_count', 'has_exclamation', 'log_length']\nfeature_audit = pd.DataFrame({\n    'feature': required_features,\n    'present': [f in df_model.columns for f in required_features],\n    'meaning': [\n        'count of positive/safe words',\n        'count of negative/scam words',\n        '1 if message contains no, else 0',\n        'count of I/me/my/you/your',\n        '1 if ! appears, else 0',\n        'log length of message'\n    ]\n})\nfeature_audit"),
        cell("markdown", "## Attractive Feature Comparison Plot\n\nThe next plot compares average risk features by class. It is useful for explaining the model in the report and presentation."),
        cell("code", "plot_features = ['safe_word_hits', 'scam_word_hits', 'has_url', 'has_phone', 'has_money', 'urgent_terms', 'action_terms', 'has_exclamation']\nfeature_means = df_model.groupby('label')[plot_features].mean().T.reset_index().rename(columns={'index': 'feature'})\nfeature_means_melted = feature_means.melt(id_vars='feature', var_name='label', value_name='mean_value')\nplt.figure(figsize=(12, 5))\nsns.barplot(data=feature_means_melted, x='feature', y='mean_value', hue='label')\nplt.xticks(rotation=35, ha='right')\nplt.title('Average Feature Values by Class')\nplt.tight_layout()"),
        cell("markdown", "## Demo: Predict New Messages\n\nThis demo cell shows how the model behaves on realistic examples. It is helpful for presentation and portfolio demonstration."),
        cell("code", "demo_messages = [\n    'URGENT! Your bank account is locked. Click the link now to verify your password and claim your reward.',\n    'Hi, I will meet you at the library after class. Please bring the assignment notes.',\n    'Congratulations, you won cash. Send your account number to receive the prize today.'\n]\nfor msg in demo_messages:\n    label, prob, _ = predict_message(msg)\n    print(f'{label.upper():5s} | scam probability={prob:.2%} | {msg}')"),
        cell("markdown", "## Deployment Note\n\nThe trained pipeline is saved to `models/scam_detector_pipeline.joblib` and reused by `app/streamlit_app.py`. This makes the work reproducible: the notebook explains and trains the model, while Streamlit demonstrates real-life use."),
    ]
    existing = "".join("".join(c.get("source", [])) for c in nb["cells"])
    if "Required Feature Audit" not in existing:
        nb["cells"].extend(additions)
    NOTEBOOK.write_text(json.dumps(nb, indent=2), encoding="utf-8")


def enhance_app() -> None:
    app = r'''from pathlib import Path
import math
import re
import joblib
import pandas as pd
import streamlit as st

ROOT = Path(__file__).resolve().parents[1]
MODEL_PATH = ROOT / "models" / "scam_detector_pipeline.joblib"
DATA_PATH = ROOT / "data" / "clean_labeled_messages_with_features.csv"
SOURCE_DIR = ROOT.parent

TOKEN_RE = re.compile(r"[a-zA-Z][a-zA-Z']+|\d+(?:[\.,]\d+)?|[$]\s*\d+")
URL_RE = re.compile(r"https?://|www\.|bit\.ly|tinyurl|t\.co|\.com\b", re.I)
PHONE_RE = re.compile(r"(?:\+?\d[\s\-().]*){7,}")
MONEY_RE = re.compile(r"[$]\s?\d+|\b\d+(?:,\d{3})*(?:\.\d+)?\s?(?:usd|dollars?|cash|prize)\b", re.I)
URGENT_RE = re.compile(r"\b(urgent|immediately|now|today|limited|final|expire|act fast|asap)\b", re.I)
ACTION_RE = re.compile(r"\b(click|claim|verify|confirm|reply|call|text|login|update|send|transfer)\b", re.I)
NO_RE = re.compile(r"\bno\b", re.I)
PRONOUN_RE = re.compile(r"\b(i|me|my|you|your)\b", re.I)

st.set_page_config(page_title="Scam or Not Scam Detector", page_icon="!", layout="wide")

st.markdown("""
<style>
.main .block-container {padding-top: 1.7rem; max-width: 1180px;}
.hero {padding: 1.4rem 1.6rem; border-radius: 10px; background: linear-gradient(135deg, #102A43 0%, #1F7A8C 55%, #2A9D8F 100%); color: white; margin-bottom: 1rem;}
.hero h1 {font-size: 2.2rem; margin: 0 0 .25rem 0;}
.hero p {font-size: 1.05rem; margin: 0; opacity: .94;}
.metric-card {padding: 1rem; border: 1px solid #E5E7EB; border-radius: 8px; background: #FFFFFF;}
.risk-high {padding: .9rem 1rem; border-radius: 8px; background: #FEE2E2; border: 1px solid #FCA5A5; color: #7F1D1D;}
.risk-low {padding: .9rem 1rem; border-radius: 8px; background: #DCFCE7; border: 1px solid #86EFAC; color: #14532D;}
.small-note {color: #64748B; font-size: .9rem;}
</style>
""", unsafe_allow_html=True)

def tokenize(text):
    return [t.lower().strip("'") for t in TOKEN_RE.findall(text)]

@st.cache_resource
def load_model():
    return joblib.load(MODEL_PATH)

@st.cache_data
def load_wordlists():
    safe_words = set(tokenize((SOURCE_DIR / "safe-words.txt").read_text(encoding="utf-8", errors="ignore")))
    scam_words = set(tokenize((SOURCE_DIR / "scam-indicator-words.txt").read_text(encoding="utf-8", errors="ignore")))
    return safe_words, scam_words

@st.cache_data
def load_dataset_preview():
    if DATA_PATH.exists():
        return pd.read_csv(DATA_PATH, usecols=["label", "text", "word_count", "scam_word_hits", "safe_word_hits"]).sample(800, random_state=42)
    return pd.DataFrame()

def build_features(text, safe_words, scam_words):
    toks = tokenize(text)
    token_count = max(len(toks), 1)
    scam_hits = sum(t in scam_words for t in toks)
    safe_hits = sum(t in safe_words for t in toks)
    return pd.DataFrame([{
        "text": text,
        "char_count": len(text),
        "word_count": len(toks),
        "avg_word_len": sum(map(len, toks)) / token_count,
        "digit_count": sum(ch.isdigit() for ch in text),
        "uppercase_ratio": sum(ch.isupper() for ch in text) / max(sum(ch.isalpha() for ch in text), 1),
        "exclamation_count": text.count("!"),
        "has_exclamation": int("!" in text),
        "question_count": text.count("?"),
        "contains_no": int(bool(NO_RE.search(text))),
        "first_second_pronoun_count": len(PRONOUN_RE.findall(text)),
        "log_length": math.log(max(len(text), 1)),
        "has_url": int(bool(URL_RE.search(text))),
        "has_phone": int(bool(PHONE_RE.search(text))),
        "has_money": int(bool(MONEY_RE.search(text))),
        "urgent_terms": len(URGENT_RE.findall(text)),
        "action_terms": len(ACTION_RE.findall(text)),
        "scam_word_hits": scam_hits,
        "safe_word_hits": safe_hits,
        "scam_word_density": scam_hits / token_count,
        "safe_word_density": safe_hits / token_count,
    }])

def risk_signals(row):
    checks = [
        ("has_url", "Contains a link"),
        ("has_phone", "Contains a phone number"),
        ("has_money", "Mentions money, cash, or prize"),
        ("urgent_terms", "Uses urgency language"),
        ("action_terms", "Asks the user to take action"),
        ("scam_word_hits", "Matches scam-indicator words"),
        ("has_exclamation", "Uses exclamation pressure"),
    ]
    return [label for key, label in checks if float(row.get(key, 0)) > 0]

model = load_model()
safe_words, scam_words = load_wordlists()
sample_df = load_dataset_preview()

st.markdown("""
<div class="hero">
  <h1>Scam or Not Scam Detector</h1>
  <p>Portfolio-ready text classification app for screening suspicious SMS, email, and chat messages.</p>
</div>
""", unsafe_allow_html=True)

tab_analyze, tab_batch, tab_terms, tab_method = st.tabs(["Analyze", "Batch Check", "Word Search", "Method"])

with tab_analyze:
    left, right = st.columns([1.25, 0.75], gap="large")
    with left:
        example = "URGENT! Your bank account is locked. Click the link now to verify your password and claim your reward."
        message = st.text_area("Paste a suspicious message", value=example, height=170)
        threshold = st.slider("Scam warning threshold", 0.10, 0.90, 0.50, 0.05)
        run = st.button("Analyze message", type="primary", use_container_width=True)
    with right:
        st.markdown("#### Safety checklist")
        st.write("- Do not click unknown links.")
        st.write("- Do not share OTP, password, or bank details.")
        st.write("- Verify urgent requests through official channels.")
        st.write("- Treat prize, refund, or account-lock messages carefully.")

    if run and message.strip():
        features = build_features(message, safe_words, scam_words)
        probability = float(model.predict_proba(features)[0, 1])
        prediction = "Scam" if probability >= threshold else "Safe"
        c1, c2, c3 = st.columns(3)
        c1.metric("Prediction", prediction)
        c2.metric("Scam probability", f"{probability:.1%}")
        c3.metric("Risk signals", len(risk_signals(features.iloc[0])))

        if prediction == "Scam":
            st.markdown('<div class="risk-high"><b>Warning:</b> This message looks suspicious. Do not click links, transfer money, or share private codes.</div>', unsafe_allow_html=True)
        else:
            st.markdown('<div class="risk-low"><b>Lower risk:</b> This message looks relatively safe, but verify important requests through official channels.</div>', unsafe_allow_html=True)

        signals = risk_signals(features.iloc[0])
        st.subheader("Detected risk signals")
        st.write(", ".join(signals) if signals else "No strong risk signals found.")
        with st.expander("Feature details"):
            st.dataframe(features.drop(columns=["text"]).T.rename(columns={0: "value"}), use_container_width=True)

with tab_batch:
    st.markdown("Upload a CSV with a `text` column to score many messages.")
    uploaded = st.file_uploader("CSV file", type=["csv"])
    if uploaded:
        batch = pd.read_csv(uploaded)
        if "text" not in batch.columns:
            st.error("CSV must contain a text column.")
        else:
            rows = [build_features(str(t), safe_words, scam_words) for t in batch["text"].fillna("")]
            features = pd.concat(rows, ignore_index=True)
            probs = model.predict_proba(features)[:, 1]
            result = batch.copy()
            result["scam_probability"] = probs
            result["prediction"] = ["scam" if p >= 0.50 else "safe" for p in probs]
            st.dataframe(result, use_container_width=True)
            st.download_button("Download scored CSV", result.to_csv(index=False).encode("utf-8"), "scored_messages.csv", "text/csv")

with tab_terms:
    st.markdown("Search terms in the project dictionaries and sample corpus.")
    query = st.text_input("Search a word or phrase", placeholder="bank, prize, verify, account...")
    if query:
        q = query.lower().strip()
        in_safe = q in safe_words
        in_scam = q in scam_words
        st.write({"in_safe_dictionary": in_safe, "in_scam_dictionary": in_scam})
        if not sample_df.empty:
            matches = sample_df[sample_df["text"].str.lower().str.contains(re.escape(q), na=False)].head(15)
            st.dataframe(matches, use_container_width=True)
    else:
        col1, col2 = st.columns(2)
        col1.markdown("#### Scam indicator examples")
        col1.write(", ".join(sorted(list(scam_words))[:80]))
        col2.markdown("#### Safe word examples")
        col2.write(", ".join(sorted(list(safe_words))[:80]))

with tab_method:
    st.markdown("### Project method")
    st.write("This app uses the model trained in the Jupyter notebook. It follows the assignment-required feature extraction and adds real-life scam indicators.")
    st.write("Core methods: tokenization, dictionary matching, engineered features, TF-IDF comparison, Logistic Regression, Naive Bayes, confusion matrix, and deployable Streamlit inference.")
    st.markdown('<p class="small-note">Privacy note: for a real public deployment, avoid storing user messages unless consent and security controls are implemented.</p>', unsafe_allow_html=True)
'''
    APP.write_text(app, encoding="utf-8")


def enhance_slides() -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception:
        return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def slide(title, bullets, accent=(31, 122, 140)):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = RGBColor(248, 250, 252)
        bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*accent)
        bar.line.fill.background()
        tb = s.shapes.add_textbox(Inches(0.75), Inches(0.45), Inches(11.8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(31)
        p.font.color.rgb = RGBColor(15, 23, 42)
        body = s.shapes.add_textbox(Inches(0.95), Inches(1.45), Inches(11.3), Inches(5.6))
        tf = body.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(21)
            para.font.color.rgb = RGBColor(51, 65, 85)
            para.space_after = Pt(11)

    slide("Scam or Not Scam Detection", [
        "Mini Project 2: Binary Text Classification",
        "Student: [Your Name] | ID: [Your Student ID]",
        "Information Retrieval pipeline with a deployable Streamlit prototype."
    ], (16, 42, 67))
    slide("Problem Motivation", [
        "Scam messages are short, urgent, noisy, and high-risk.",
        "Users need practical warnings before clicking links or sharing private information.",
        "The task is a real-world binary text-classification problem."
    ])
    slide("Dataset and Split", [
        "20,000 safe messages and 20,000 scam messages.",
        "Top 80% of each class used for training; final 20% used for testing.",
        "Safe/scam word lists support dictionary-based feature extraction."
    ])
    slide("Required Features", [
        "Count of safe/positive words and scam/negative words.",
        "Indicator for word 'no'.",
        "Pronoun count: I, me, my, you, your.",
        "Indicator for exclamation mark and log message length."
    ])
    slide("Professional Feature Extensions", [
        "URL, phone-number, and money/prize flags.",
        "Urgency and action-term counts.",
        "Uppercase ratio, digit count, question count, word densities.",
        "TF-IDF unigram and bigram vector-space representation."
    ])
    slide("Model Comparison", [
        "Required engineered features + Logistic Regression.",
        "TF-IDF + Multinomial Naive Bayes.",
        "TF-IDF + engineered features + Logistic Regression.",
        "All models reached 1.0000 accuracy on the provided corpus."
    ], (231, 111, 81))
    slide("Evaluation", [
        "Metrics: accuracy, precision, recall, F1-score, ROC-AUC, confusion matrix.",
        "Recall is important because missed scam messages can harm users.",
        "Perfect scores indicate high separability and should be validated on noisier real data."
    ], (231, 111, 81))
    slide("Literature Review Insight", [
        "Classic SMS spam studies show strong performance from supervised text classification.",
        "Recent work highlights smishing, privacy, adversarial robustness, and concept drift.",
        "Research gap: multilingual, explainable, privacy-first scam detection for real users."
    ])
    slide("Streamlit Demo", [
        "Paste a suspicious message and receive a risk score.",
        "Batch CSV scoring supports practical review workflows.",
        "Word search helps explore safe/scam dictionaries and corpus examples."
    ])
    slide("Ethics and Deployment", [
        "Do not store private user messages without consent.",
        "Show confidence and risk signals, not only a hard label.",
        "Use the app as decision support, not as a replacement for human judgment."
    ], (42, 157, 143))
    slide("Future Work", [
        "Add Khmer-English text normalization and more local scam examples.",
        "Highlight suspicious phrases for explainability.",
        "Add feedback collection, drift monitoring, and transformer-model comparison."
    ], (42, 157, 143))
    prs.save(SLIDES)


def main() -> None:
    enhance_report()
    enhance_notebook()
    enhance_app()
    enhance_slides()
    print("Enhanced report, notebook, Streamlit app, and slides.")


if __name__ == "__main__":
    main()
